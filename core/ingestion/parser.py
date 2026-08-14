"""
Universal document parser - supports PDF, Excel, PPT, Word, Images, Markdown, HTML, TXT
Fully generalized, no industry-specific content
"""
import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# Attempt to import parsing libraries
try:
    import pandas as pd
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    from pptx import Presentation
    HAS_PPT = True
except ImportError:
    HAS_PPT = False

try:
    import pytesseract
    from PIL import Image
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

try:
    from bs4 import BeautifulSoup
    HAS_BS4 = True
except ImportError:
    HAS_BS4 = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pdf2image import convert_from_path
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False

try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

from ..config import GRID_RECONSTRUCTION_ENABLED, INGEST_MAX_WORKERS, settings
from ..models import get_model_client


class ParsedPage:
    """Parsed page"""
    def __init__(self, page_num: int, raw_text: str = "",
                 section_title: str = "", content_dict: dict = None,
                 page_image: Any = None):
        self.page_num = page_num
        self.raw_text = raw_text
        self.section_title = section_title
        self.content_dict = content_dict or {}
        self.page_image = page_image

    def to_dict(self) -> dict:
        return {
            "page_num": self.page_num,
            "raw_text": self.raw_text,
            "section_title": self.section_title,
            "content_dict": self.content_dict,
        }


class ParsedDocument:
    """Parsed document"""
    def __init__(self):
        self.filename = ""
        self.original_path = ""
        self.file_size = 0
        self.pages: list[ParsedPage] = []
        self.metadata: dict[str, Any] = {}

    @property
    def total_pages(self):
        return len(self.pages)

    def to_dict(self) -> dict:
        return {
            "filename": self.filename,
            "original_path": self.original_path,
            "file_size": self.file_size,
            "total_pages": self.total_pages,
            "metadata": self.metadata,
            "pages": [p.to_dict() for p in self.pages],
        }


class DocumentParser:
    """Universal document parser - fully generalized, no industry-specific logic"""

    SUPPORTED_EXTENSIONS = {
        ".pdf", ".md", ".txt", ".html", ".htm",
        ".xlsx", ".xls", ".pptx", ".ppt",
        ".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".gif",
        ".docx"
    }

    def __init__(self):
        self._page_classifier = None  # Lazy-load PageClassifier

    @property
    def page_classifier(self):
        """Lazy-load PageClassifier to avoid circular imports"""
        if self._page_classifier is None:
            from .layout.page_classifier import PageClassifier
            self._page_classifier = PageClassifier()
        return self._page_classifier

    def parse(self, file_path: str) -> ParsedDocument:
        """Parse document, return ParsedDocument object"""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = path.suffix.lower()
        if suffix not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file format: {suffix}")

        if suffix == ".pdf":
            return self._parse_pdf(path)
        elif suffix in [".xlsx", ".xls"]:
            return self._parse_excel(path)
        elif suffix in [".pptx", ".ppt"]:
            return self._parse_ppt(path)
        elif suffix in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".gif"]:
            return self._parse_image(path)
        elif suffix == ".md":
            return self._parse_markdown(path)
        elif suffix == ".txt":
            return self._parse_text(path)
        elif suffix in [".html", ".htm"]:
            return self._parse_html(path)
        elif suffix == ".docx":
            return self._parse_docx(path)
        else:
            raise ValueError(f"Unsupported file format: {suffix}")

    # ========================================================================
    # PDF Parsing
    # ========================================================================

    def _parse_pdf(self, path: Path) -> ParsedDocument:
        """Parse PDF file - pypdf for metadata/TOC, pdfplumber for text/tables, pdf2image for rendering"""
        doc = ParsedDocument()
        doc.filename = path.name
        doc.original_path = str(path.absolute())
        doc.file_size = path.stat().st_size

        if not HAS_PYPDF or not HAS_PDFPLUMBER:
            logger.error("pypdf or pdfplumber unavailable, cannot parse PDF")
            doc.pages.append(ParsedPage(page_num=1, raw_text="PDF parsing failed: pypdf/pdfplumber not installed"))
            return doc

        vlm_analyses = {}

        # ── Two-pass VLM classification (pre-filter before VLM) ──
        # Pass 1: Fast image detection via pypdf + text pre-filter via pdfplumber
        # Only pages with images AND minimal text go to VLM
        page_classes = {}
        vlm_candidate_pages = []
        total_pages = 0

        try:
            pypdf_reader = pypdf.PdfReader(str(path))
            plumber_doc = pdfplumber.open(str(path))
            total_pages = len(pypdf_reader.pages)

            for page_num in range(total_pages):
                pdf_page_num = page_num + 1
                has_images = False

                # Pass 1a: Check for images via pypdf /Resources/XObject
                try:
                    page_obj = pypdf_reader.pages[page_num]
                    resources = page_obj.get("/Resources")
                    if resources:
                        xobjects = resources.get("/XObject")
                        if xobjects:
                            for obj in xobjects.values():
                                obj_ref = obj.get_object() if hasattr(obj, "get_object") else obj
                                if hasattr(obj_ref, "get") and obj_ref.get("/Subtype") == "/Image":
                                    has_images = True
                                    break
                except Exception:
                    pass

                # Pass 1b: Extract text length via pdfplumber to decide whether VLM is needed
                try:
                    plumber_page = plumber_doc.pages[page_num]
                    text = plumber_page.extract_text() or ""
                except Exception:
                    text = ""
                text = text.replace('\x00', '')
                text_len = len(text.strip())

                vlm_min_text = settings.CHART_CONFIG.get("vlm_min_text_len_for_candidate", 1000)

                if not has_images:
                    # No images -> definitely not a visual candidate.
                    # If the page is also empty text, mark it as BLANK so downstream
                    # chunking and retrieval can skip it without any VLM work.
                    if text_len == 0:
                        page_classes[pdf_page_num] = "BLANK"
                    else:
                        page_classes[pdf_page_num] = "TEXT"
                    continue

                if text_len > vlm_min_text:
                    # Lots of text -> images are likely decorative/icons, skip VLM
                    page_classes[pdf_page_num] = "TEXT"
                    continue

                # Pass 1c: Has images + minimal text -> VLM candidate
                vlm_candidate_pages.append(pdf_page_num)

            plumber_doc.close()
        except Exception as e:
            logger.warning(f"PDF pre-filter failed, falling back to all-TEXT: {e}")
            # Fallback: mark all as TEXT, skip VLM (no blank detection on failure path)
            try:
                pypdf_reader = pypdf.PdfReader(str(path))
                total_pages = len(pypdf_reader.pages)
                for pn in range(1, total_pages + 1):
                    page_classes[pn] = "TEXT"
            except Exception:
                pass

        # Pass 2: VLM classification only for candidate pages (with images + minimal text)
        page_images = {}
        if vlm_candidate_pages:
            page_images = self._render_pdf_pages(str(path), dpi=72)
            candidate_images = {
                pn: page_images[pn] for pn in vlm_candidate_pages if pn in page_images
            }
            if candidate_images:
                from concurrent.futures import ThreadPoolExecutor, as_completed
                max_workers = max(1, min(INGEST_MAX_WORKERS, len(candidate_images)))
                logger.info(
                    f"VLM classification: {len(candidate_images)} candidate pages "
                    f"(of {total_pages} total) with {max_workers} workers"
                )
                with ThreadPoolExecutor(max_workers=max_workers) as executor:
                    futures = {
                        executor.submit(self._classify_one_page, img, pn): pn
                        for pn, img in candidate_images.items()
                    }
                    for future in as_completed(futures):
                        pn = futures[future]
                        try:
                            page_classes[pn] = future.result(timeout=30)
                        except Exception as e:
                            logger.warning(f"Page {pn} classification failed: {e}")
                            page_classes[pn] = "TEXT"
                stats = self.page_classifier.stats
                logger.info(
                    f"VLM classification complete: {stats['total']} pages "
                    f"(chart={stats['chart']}, image={stats['image']}, text={stats['text']}, blank={stats['blank']})"
                )
                self.page_classifier.reset_stats()
            else:
                # No renderable candidate pages -> all TEXT
                for pn in vlm_candidate_pages:
                    page_classes[pn] = "TEXT"
        else:
            logger.info(f"VLM classification: 0 candidate pages (of {total_pages} total) -> all TEXT")

        # Ensure all pages have a class (fallback to TEXT)
        for pn in range(1, total_pages + 1):
            if pn not in page_classes:
                page_classes[pn] = "TEXT"

        try:
            pypdf_reader = pypdf.PdfReader(str(path))
            plumber_doc = pdfplumber.open(str(path))
            total_pages = len(pypdf_reader.pages)
            doc.metadata["num_pages"] = total_pages

            # Extract TOC from pypdf outline
            outline_map = self._build_outline_map(pypdf_reader)
            if outline_map:
                logger.info(f"PDF bookmarks extracted: {len(outline_map)} pages have section titles")
            # FIX: Build toc from ALL outline entries (not page-deduped outline_map).
            # Datasheets pack multiple small sections on one page (e.g. 1.2.4 Video Codec
            # and 1.2.5 Neural Process Unit both on p12); page-keyed dedup silently
            # dropped those sections from the structure index.
            full_toc = self._build_full_toc(pypdf_reader)
            if full_toc:
                doc.metadata["toc"] = full_toc

            for page_num in range(total_pages):
                pdf_page_num = page_num + 1

                # 1. Extract text + tables via pdfplumber
                plumber_page = None
                try:
                    plumber_page = plumber_doc.pages[page_num]
                    text = plumber_page.extract_text() or ""
                except Exception:
                    text = ""
                text = text.replace('\x00', '')

                # 2. Extract tables via pdfplumber
                table_md = ""
                if plumber_page is not None:
                    try:
                        raw_tables = plumber_page.extract_tables()
                        if raw_tables:
                            for tab_idx, tab in enumerate(raw_tables):
                                if not tab or len(tab) < 2:
                                    continue
                                try:
                                    import pandas as pd
                                    df = pd.DataFrame(tab[1:], columns=tab[0])
                                    if df is not None and not df.empty:
                                        # Filter pseudo-tables
                                        total_cells = df.shape[0] * df.shape[1]
                                        empty_cells = df.isna().sum().sum() + (df == '').sum().sum()
                                        empty_ratio = empty_cells / total_cells if total_cells > 0 else 0
                                        if empty_ratio >= 0.5 and df.shape[0] <= 3 and df.shape[1] <= 3:
                                            logger.debug(f"Skipping pseudo-table page {pdf_page_num} table {tab_idx+1}")
                                            continue
                                        md = df.to_markdown(index=False, floatfmt='')
                                        table_md += f"\n\n[Table {tab_idx+1}]\n{md}\n"
                                except Exception as e:
                                    logger.debug(f"Table to markdown failed page {pdf_page_num} table {tab_idx+1}: {e}")
                    except Exception as e:
                        logger.debug(f"Table detection failed page {pdf_page_num}: {e}")

                full_text = text
                if table_md:
                    full_text += table_md

                # 2b. Rebuild labeled ruled grids (e.g. ball maps) from vector info.
                # Text extraction scrambles grid diagrams into misleading garbage;
                # the reconstruction is deterministic and model-free.
                if GRID_RECONSTRUCTION_ENABLED and plumber_page is not None:
                    try:
                        from .layout.grid_reconstructor import reconstruct_grid_table
                        grid_md = reconstruct_grid_table(plumber_page)
                        if grid_md:
                            full_text += f"\n\n[Grid Map]\n{grid_md}\n"
                    except Exception as e:
                        logger.debug(f"Grid reconstruction failed page {pdf_page_num}: {e}")

                # 3. VLM page classification (pre-computed in two-pass above)
                page_image = page_images.get(pdf_page_num)
                vlm_needed = False

                page_class = page_classes.get(pdf_page_num, "TEXT")

                if page_class == "BLANK":
                    # Blank page: no visual content, keep text empty and do not run VLM
                    logger.info(f"PDF p{pdf_page_num}: page_class=BLANK -> skip VLM and chunking")
                    full_text = ""

                elif page_class == "CHART":
                    # Chart page -> high DPI VLM deep analysis
                    logger.info(f"PDF p{pdf_page_num}: page_class=CHART -> high DPI deep analysis")
                    hi_res_img = self._render_single_page(str(path), pdf_page_num, dpi=150)
                    if hi_res_img:
                        analysis = self._analyze_pdf_page_with_vlm(
                            hi_res_img, pdf_page_num, page_text=full_text
                        )
                        if analysis:
                            vlm_analyses[pdf_page_num] = analysis
                            full_text += (
                                f"\n\n---\n\n### Page Visual Analysis (VLM)\n\n{analysis}\n"
                            )
                            vlm_needed = True

                elif page_class == "IMAGE":
                    # Image / scanned page -> generic VLM description/transcription
                    logger.info(f"PDF p{pdf_page_num}: page_class=IMAGE -> VLM image description")
                    img_cfg = settings.CHART_CONFIG
                    if img_cfg.get("vlm_image_description_enabled", True):
                        # Reuse the already rendered 72 DPI image to avoid a second render
                        # unless a higher-resolution image is requested for description.
                        desc_img = page_image
                        if img_cfg.get("vlm_image_description_use_hi_res", False):
                            desc_img = self._render_single_page(str(path), pdf_page_num, dpi=150) or page_image
                        if desc_img:
                            analysis = self._describe_pdf_page_with_vlm(
                                desc_img, pdf_page_num, page_text=full_text
                            )
                            if analysis:
                                vlm_analyses[pdf_page_num] = analysis
                                full_text += (
                                    f"\n\n---\n\n### Page Visual Analysis (VLM)\n\n{analysis}\n"
                                )
                                vlm_needed = True
                    else:
                        logger.info(f"PDF p{pdf_page_num}: image description disabled by config")

                else:
                    logger.debug(f"PDF p{pdf_page_num}: page_class=TEXT -> text extraction")

                # First page metadata
                if pdf_page_num == 1:
                    meta = pypdf_reader.metadata
                    if meta:
                        doc.metadata.update({
                            "title": getattr(meta, "title", None) or (meta.get("/Title") if isinstance(meta, dict) else None),
                            "author": getattr(meta, "author", None) or (meta.get("/Author") if isinstance(meta, dict) else None),
                            "subject": getattr(meta, "subject", None) or (meta.get("/Subject") if isinstance(meta, dict) else None),
                        })

                parsed_page = ParsedPage(
                    page_num=pdf_page_num,
                    raw_text=full_text,
                    section_title=self._extract_section_title(full_text, pdf_page_num, outline_map),
                    page_image=page_image,
                    content_dict={
                        "pdf_text": text,
                        "has_tables": bool(table_md),
                        "vlm_analysis": vlm_analyses.get(pdf_page_num),
                        "page_class": page_class,
                        "vlm_needed": vlm_needed,
                    }
                )
                doc.pages.append(parsed_page)

            plumber_doc.close()

        except Exception as e:
            logger.error(f"PDF parsing failed {path}: {e}")
            if not doc.pages:
                # pdfplumber + pypdf failed entirely.
                # Fallback: extract text page-by-page with pymupdf (MuPDF),
                # which is more lenient with corrupted PDF content streams.
                # Page boundaries must be preserved: merging all pages into a
                # single record disables page-level retrieval and structure
                # indexing for exactly the documents that need this fallback.
                page_texts = self._extract_pages_with_pymupdf(str(path))
                if page_texts:
                    for i, text in enumerate(page_texts, start=1):
                        doc.pages.append(ParsedPage(
                            page_num=i,
                            raw_text=f"[Partial extraction via MuPDF fallback]\n\n{text}",
                            content_dict={"pdf_text": text, "fallback": "pymupdf"},
                        ))
                else:
                    doc.pages.append(ParsedPage(page_num=1, raw_text=f"PDF parsing failed: {e}"))

        return doc

    @staticmethod
    def _extract_pages_with_pymupdf(path: str) -> list[str]:
        """Fallback extraction using pymupdf (MuPDF).

        More lenient than pdfplumber/pypdf for PDFs with corrupted content
        streams. Returns one text entry per non-empty page, preserving page
        order so the caller can keep page-level addressing intact.
        """
        try:
            import fitz  # pymupdf
            doc = fitz.open(path)
            page_texts = []
            for i in range(doc.page_count):
                try:
                    t = doc[i].get_text()
                    if t and t.strip():
                        page_texts.append(t)
                except Exception:
                    pass
            doc.close()
            return page_texts
        except Exception as e:
            logger.warning(f"[PARSER] pymupdf fallback also failed: {e}")
            return []

    def _classify_one_page(self, page_image, page_num: int) -> str:
        """Single-page VLM classification wrapper (for concurrent use)"""
        try:
            return self.page_classifier.classify(page_image, page_num)
        except Exception as e:
            logger.warning(f"Classify p{page_num} failed: {e}")
            return "TEXT"

    def _classify_pdf_page(self, page, text: str, has_images: bool, table_md: str) -> str:
        """PDF page type classification

        Return types:
        - "text": plain text page (no images, no tables)
        - "native_table": native text table page (no images, has text tables)
        - "image_table": embedded image table page (has images, minimal text)
        - "scan_page": scanned/complex image page (has images, almost no text)
        """
        text_len = len(text.strip())
        has_table_markers = bool(table_md) or "Table" in text or "Fig." in text or "Figure" in text

        # 1. Scanned/complex image page: has images + very little text (<100 chars)
        if has_images and text_len < 100:
            return "scan_page"

        # 2. Embedded image table page: has images + minimal text (<500 chars) + incomplete text
        # Check text completeness: if text has section headings and body text, consider it complete
        text_is_complete = (
            len(text.strip().split('\n')) > 5  # at least 5 lines
            and any(c in text for c in ['章', '节', '1.', '2.', '3.'])  # has chapter structure
        )
        if has_images and text_len < 500 and not text_is_complete:
            return "image_table"

        # 3. Native text table page: no images + has table markers + has text
        if not has_images and has_table_markers and text_len > 50:
            return "native_table"

        # 4. Plain text page: no images + no table markers
        if not has_images and not has_table_markers:
            return "text"

        # 5. Default: plain text page
        return "text"

    def _render_single_page(self, pdf_path: str, page_num: int, dpi: int = 150):
        """Render a single PDF page as PIL Image via pdf2image"""
        if HAS_PDF2IMAGE:
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(
                    pdf_path, dpi=dpi,
                    first_page=page_num, last_page=page_num
                )
                if images:
                    return images[0]
            except Exception as e:
                logger.warning(f"pdf2image single page render failed p{page_num}: {e}")

        return None

    def _render_pdf_pages(self, pdf_path: str, dpi: int = 100) -> dict[int, Any]:
        """Render PDF pages to PIL Images via pdf2image"""
        page_images = {}
        if HAS_PDF2IMAGE:
            try:
                images = convert_from_path(pdf_path, dpi=dpi)
                for i, img in enumerate(images, 1):
                    page_images[i] = img
                logger.info(f"pdf2image rendered {len(images)} pages")
                return page_images
            except Exception as e:
                logger.warning(f"pdf2image failed: {e}")

        return page_images

    def _analyze_pdf_page_with_vlm(self, page_image, page_num: int, page_text: str = "") -> str:
        """Analyze PDF page with VLM"""
        try:
            import os
            import tempfile

            client = get_model_client()

            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                page_image.save(tmp.name, 'PNG')
                tmp_path = tmp.name

            try:
                text_section = ""
                if page_text and len(page_text.strip()) > 10:
                    text_section = f"""
Raw content extracted from this page (for reference):
---
{page_text[:1500]}
---
"""

                prompt = f"""You are analyzing a screenshot of a PDF page.{text_section}
Based on the page content type, focus on extracting core information:
- If it's an architecture/flow diagram: describe system module composition, interfaces between modules, data flow.
- If it's a data chart/table: extract key data points, trends, comparisons.
- If it's a product/object image: describe product appearance, interfaces, markings.
- If it's a screenshot/diagram: describe the core information.

Output in Markdown format, including:
1. **Page Type**
2. **Key Content Description**
3. **Text Annotation Transcription**
4. **Core Information Extraction**

Please output in English."""

                analysis = client.generate_with_image(
                    prompt=prompt,
                    image_path=tmp_path,
                    max_tokens=2048,
                    temperature=0.2
                )
                return analysis.strip() if analysis else ""
            finally:
                os.unlink(tmp_path)
        except Exception as e:
            logger.warning(f"PDF page {page_num} VLM analysis failed: {e}")
            return ""

    def _describe_pdf_page_with_vlm(self, page_image, page_num: int, page_text: str = "") -> str:
        """Generate a generic textual description/transcription for an image-heavy page.

        This is intentionally more general than the chart analyzer: it handles photographs,
        scanned pages, screenshots, and any non-chart visual page without domain-specific
        assumptions. The description is appended to the page's raw_text so it becomes
        searchable.
        """
        try:
            import os
            import tempfile

            client = get_model_client()
            img_cfg = settings.CHART_CONFIG

            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                page_image.save(tmp.name, 'PNG')
                tmp_path = tmp.name

            try:
                text_section = ""
                if page_text and len(page_text.strip()) > 10:
                    text_section = f"""\nRaw content already extracted from this page (use as reference only):\n---\n{page_text[:1000]}\n---\n"""

                prompt = f"""You are analyzing a document page image.{text_section}
Describe the page content accurately and concisely:
- If the page contains readable text, transcribe it.
- If the page contains a diagram, illustration, or drawing, describe the elements, labels, structure, and relationships shown.
- If the page contains a photograph or screenshot, describe what is visible and any readable text or markings.

Output in plain Markdown. Be factual and avoid guessing information not visible in the image."""

                max_tokens = img_cfg.get("vlm_image_description_max_tokens", 1024)
                temperature = img_cfg.get("vlm_image_description_temperature", 0.2)
                analysis = client.generate_with_image(
                    prompt=prompt,
                    image_path=tmp_path,
                    max_tokens=max_tokens,
                    temperature=temperature
                )
                return analysis.strip() if analysis else ""
            finally:
                os.unlink(tmp_path)
        except Exception as e:
            logger.warning(f"PDF page {page_num} VLM image description failed: {e}")
            return ""

    # ========================================================================
    # Excel Parsing
    # ========================================================================

    def _parse_excel(self, path: Path) -> ParsedDocument:
        doc = ParsedDocument()
        doc.filename = path.name
        doc.original_path = str(path.absolute())
        doc.file_size = path.stat().st_size

        if HAS_EXCEL:
            try:
                xl = pd.ExcelFile(path)
                doc.metadata["sheets"] = xl.sheet_names

                for sheet_idx, sheet_name in enumerate(xl.sheet_names, 1):
                    df = pd.read_excel(path, sheet_name=sheet_name)
                    sheet_text = f"=== Sheet: {sheet_name} ===\n"
                    sheet_text += " | ".join(map(str, df.columns)) + "\n"
                    sheet_text += "-|-".join(["---"] * len(df.columns)) + "\n"
                    for _, row in df.iterrows():
                        sheet_text += " | ".join(map(str, row.values)) + "\n"

                    parsed_page = ParsedPage(
                        page_num=sheet_idx,
                        raw_text=sheet_text,
                        section_title=sheet_name,
                        content_dict={"sheet": sheet_name, "columns": list(df.columns)}
                    )
                    doc.pages.append(parsed_page)
            except Exception as e:
                logger.error(f"Excel parsing failed {path}: {e}")

        return doc

    # ========================================================================
    # PPT Parsing
    # ========================================================================

    def _parse_ppt(self, path: Path) -> ParsedDocument:
        doc = ParsedDocument()
        doc.filename = path.name
        doc.original_path = str(path.absolute())
        doc.file_size = path.stat().st_size

        if HAS_PPT:
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                prs = Presentation(path)
                doc.metadata["num_slides"] = len(prs.slides)

                for slide_idx, slide in enumerate(prs.slides, 1):
                    slide_text = f"=== Slide {slide_idx} ===\n"
                    chart_texts = []
                    image_paths = []

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text += shape.text + "\n"

                        if getattr(shape, "has_chart", False):
                            try:
                                chart_text = self._extract_ppt_chart_data(shape.chart, slide_idx)
                                if chart_text:
                                    chart_texts.append(chart_text)
                            except Exception as e:
                                logger.debug(f"PPT chart extraction failed: {e}")

                        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image = shape.image
                                img_filename = f"ppt_{path.stem}_s{slide_idx}_{len(image_paths)}.{image.ext}"
                                img_path = settings.UPLOAD_DIR / img_filename
                                with open(img_path, "wb") as f:
                                    f.write(image.blob)
                                desc = self._describe_image(str(img_path))
                                if desc:
                                    image_paths.append({"path": str(img_path), "description": desc})
                            except Exception as e:
                                logger.debug(f"PPT image extraction failed: {e}")

                    if chart_texts:
                        slide_text += "\n[Chart data in slide]\n"
                        slide_text += "\n".join(chart_texts) + "\n"
                    if image_paths:
                        slide_text += "\n[Image content in slide]\n"
                        for img_info in image_paths:
                            slide_text += f"Image description: {img_info['description']}\n"

                    content_dict = {}
                    if image_paths:
                        content_dict["ppt_images"] = image_paths
                    if chart_texts:
                        content_dict["ppt_charts"] = chart_texts

                    parsed_page = ParsedPage(
                        page_num=slide_idx,
                        raw_text=slide_text,
                        section_title=f"Slide {slide_idx}",
                        content_dict=content_dict
                    )
                    doc.pages.append(parsed_page)
            except Exception as e:
                logger.error(f"PPT parsing failed {path}: {e}")

        return doc

    def _extract_ppt_chart_data(self, chart, slide_idx: int) -> str:
        try:
            lines = [f"Chart (Slide {slide_idx}):"]
            if hasattr(chart, "chart_title") and chart.chart_title:
                title = chart.chart_title.text_frame.text if chart.chart_title.has_text_frame else ""
                if title:
                    lines.append(f"Title: {title}")

            categories = []
            if hasattr(chart, "plots") and chart.plots:
                for plot in chart.plots:
                    if hasattr(plot, "categories"):
                        categories = [cat.label if hasattr(cat, "label") else str(cat) for cat in plot.categories]

            for series in chart.series:
                series_name = series.name if series.name else "Unnamed series"
                lines.append(f"Series: {series_name}")
                values = []
                for idx, point in enumerate(series.points):
                    cat = categories[idx] if idx < len(categories) else f"Item {idx+1}"
                    val = point.value if hasattr(point, "value") else "N/A"
                    values.append(f"  {cat}: {val}")
                lines.extend(values[:30])
                if len(series.points) > 30:
                    lines.append(f"  ... ({len(series.points)} data points total)")

            return "\n".join(lines)
        except Exception as e:
            logger.debug(f"Chart data extraction failed: {e}")
            return ""

    # ========================================================================
    # Image Parsing
    # ========================================================================

    def _parse_image(self, path: Path) -> ParsedDocument:
        doc = ParsedDocument()
        doc.filename = path.name
        doc.original_path = str(path.absolute())
        doc.file_size = path.stat().st_size

        # Extract EXIF
        exif_text = self._extract_exif(path)

        # VLM Analysis
        image_text = ""
        try:
            client = get_model_client()
            prompt = """Please describe the content of this image in detail. If the image contains text, tables, product models, technical specifications, etc., transcribe them as completely as possible. Please output in English."""
            image_text = client.generate_with_image(
                prompt=prompt,
                image_path=str(path),
                max_tokens=2048,
                temperature=0.2
            )
            if image_text and image_text.strip():
                logger.info(f"LLM image parsing succeeded: {path.name} ({len(image_text.strip())} chars)")
            else:
                image_text = ""
        except Exception as e:
            logger.warning(f"Multimodal LLM image parsing failed {path}: {e}")

        # OCR fallback
        if not image_text:
            image_text = self._ocr_image_fallback(path)
            if image_text.startswith("Image file:"):
                image_text = ""

        parts = []
        if exif_text:
            parts.append(exif_text)
        if image_text and image_text.strip():
            parts.append(image_text.strip())

        text = "\n\n".join(parts) if parts else f"Image file: {path.name}"

        doc.pages.append(ParsedPage(page_num=1, raw_text=text))
        return doc

    def _extract_exif(self, path: Path) -> str:
        try:
            from PIL.ExifTags import GPSTAGS, TAGS
            with Image.open(path) as img:
                exif = img._getexif()
                if not exif:
                    return ""
                info = {}
                gps_info = {}
                for tag_id, value in exif.items():
                    tag_name = TAGS.get(tag_id, tag_id)
                    if tag_name == 'GPSInfo':
                        for gps_tag_id, gps_value in value.items():
                            gps_tag_name = GPSTAGS.get(gps_tag_id, gps_tag_id)
                            gps_info[gps_tag_name] = gps_value
                    else:
                        info[tag_name] = value

                parts = []
                time_str = info.get('DateTimeOriginal') or info.get('DateTime')
                if time_str:
                    parts.append(f"Capture time: {time_str}")

                lat = gps_info.get('GPSLatitude')
                lat_ref = gps_info.get('GPSLatitudeRef')
                lon = gps_info.get('GPSLongitude')
                lon_ref = gps_info.get('GPSLongitudeRef')
                if lat and lat_ref and lon and lon_ref:
                    lat_str = self._convert_gps_coords(lat, lat_ref)
                    lon_str = self._convert_gps_coords(lon, lon_ref)
                    parts.append(f"Capture location: Lat {lat_str}, Lon {lon_str}")

                return "\n".join(parts) if parts else ""
        except Exception as e:
            logger.debug(f"EXIF extraction failed {path}: {e}")
            return ""

    def _convert_gps_coords(self, coords, ref):
        try:
            if isinstance(coords, tuple) and len(coords) == 3:
                def to_float(v):
                    if hasattr(v, 'numerator') and hasattr(v, 'denominator'):
                        return float(v.numerator) / float(v.denominator)
                    return float(v)
                degrees = to_float(coords[0])
                minutes = to_float(coords[1])
                seconds = to_float(coords[2])
                decimal = degrees + minutes / 60 + seconds / 3600
                if ref in ('S', 'W'):
                    decimal = -decimal
                return f"{decimal:.6f}°"
            return str(coords)
        except Exception:
            return str(coords)

    def _ocr_image_fallback(self, path: Path) -> str:
        text = f"Image file: {path.name}"
        if HAS_OCR:
            try:
                image = Image.open(path)
                ocr_text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                if ocr_text and ocr_text.strip():
                    text = ocr_text.strip()
                    logger.info(f"OCR parsing succeeded: {path.name}")
            except Exception as e:
                logger.error(f"OCR failed {path}: {e}")
        return text

    def _describe_image(self, image_path: str) -> str:
        try:
            client = get_model_client()
            prompt = "Please describe the content of this image in detail. Please output in English."
            desc = client.generate_with_image(
                prompt=prompt,
                image_path=image_path,
                max_tokens=512,
                temperature=0.2
            )
            return desc.strip() if desc else ""
        except Exception as e:
            logger.debug(f"Image description failed: {e}")
            return ""

    # ========================================================================
    # Markdown / Text / HTML / DOCX
    # ========================================================================

    def _parse_markdown(self, path: Path) -> ParsedDocument:
        doc = ParsedDocument()
        doc.filename = path.name
        doc.original_path = str(path.absolute())
        doc.file_size = path.stat().st_size

        with open(path, encoding='utf-8') as f:
            content = f.read()

        chunk_texts = self._simple_chunk_by_heading(content)
        for i, chunk in enumerate(chunk_texts, 1):
            doc.pages.append(ParsedPage(
                page_num=i,
                raw_text=chunk,
                section_title=self._extract_section_title(chunk, i)
            ))
        return doc

    def _parse_text(self, path: Path) -> ParsedDocument:
        doc = ParsedDocument()
        doc.filename = path.name
        doc.original_path = str(path.absolute())
        doc.file_size = path.stat().st_size

        with open(path, encoding='utf-8') as f:
            content = f.read()

        chunk_texts = self._chunk_by_structure(content, max_chunk_size=2000)
        if len(chunk_texts) <= 1 and len(content) > 2000:
            chunk_texts = self._simple_chunk(content, chunk_size=2000)

        for i, chunk in enumerate(chunk_texts, 1):
            doc.pages.append(ParsedPage(
                page_num=i,
                raw_text=chunk,
                section_title=self._extract_section_title(chunk, i)
            ))
        return doc

    def _parse_html(self, path: Path) -> ParsedDocument:
        doc = ParsedDocument()
        doc.filename = path.name
        doc.original_path = str(path.absolute())
        doc.file_size = path.stat().st_size

        text = ""
        if HAS_BS4:
            with open(path, encoding='utf-8') as f:
                content = f.read()
            soup = BeautifulSoup(content, 'html.parser')
            text = soup.get_text()

        chunk_texts = self._chunk_by_structure(text, max_chunk_size=2000)
        if len(chunk_texts) <= 1 and len(text) > 2000:
            chunk_texts = self._simple_chunk(text, chunk_size=2000)

        for i, chunk in enumerate(chunk_texts, 1):
            doc.pages.append(ParsedPage(
                page_num=i,
                raw_text=chunk,
                section_title=self._extract_section_title(chunk, i)
            ))
        return doc

    def _parse_docx(self, path: Path) -> ParsedDocument:
        if not HAS_DOCX:
            raise ImportError("python-docx not installed")

        doc = ParsedDocument()
        doc.filename = path.name
        doc.original_path = str(path.absolute())
        doc.file_size = path.stat().st_size

        d = DocxDocument(str(path))

        paragraphs = []
        for para in d.paragraphs:
            text = para.text.strip()
            if text:
                style = para.style.name if para.style else ''
                if 'Heading' in style or 'heading' in style or '标题' in style:
                    level = 1
                    for c in style:
                        if c.isdigit():
                            level = int(c)
                            break
                    paragraphs.append('#' * level + ' ' + text)
                else:
                    paragraphs.append(text)

        table_texts = []
        for table in d.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                rows.append(' | '.join(cells))
            if rows:
                header = rows[0]
                separator = '|'.join(['---' for _ in rows[0].split('|')])
                table_texts.append(f'| {header} |\n|{separator}|\n' + '\n'.join(f'| {r} |' for r in rows[1:]))

        full_text = '\n\n'.join(paragraphs)
        if table_texts:
            full_text += '\n\n' + '\n\n'.join(table_texts)

        chunk_texts = self._chunk_by_structure(full_text, max_chunk_size=2000)
        if len(chunk_texts) <= 1 and len(full_text) > 2000:
            chunk_texts = self._simple_chunk(full_text, chunk_size=2000)

        for i, chunk in enumerate(chunk_texts, 1):
            doc.pages.append(ParsedPage(
                page_num=i,
                raw_text=chunk,
                section_title=self._extract_section_title(chunk, i)
            ))
        return doc

    # ========================================================================
    # Text Chunking Utilities
    # ========================================================================

    def _simple_chunk(self, text: str, chunk_size: int = 2000, overlap: int = 100) -> list[str]:
        if len(text) <= chunk_size:
            return [text]
        heading_chunks = self._chunk_by_structure(text, max_chunk_size=chunk_size)
        if heading_chunks and len(heading_chunks) > 1:
            return heading_chunks

        chunks = []
        start = 0
        text_length = len(text)
        while start < text_length:
            end = min(start + chunk_size, text_length)
            chunks.append(text[start:end])
            start = end
            if start >= text_length:
                break
            start = max(0, start - overlap)
            if start >= end:
                break
        return chunks

    def _chunk_by_structure(self, text: str, max_chunk_size: int = 2000) -> list[str]:
        import re
        lines = text.split("\n")
        chunks = []
        current_chunk = []
        current_size = 0

        def _flush_chunk():
            nonlocal current_chunk, current_size
            if current_chunk:
                chunk_text = "\n".join(current_chunk).strip()
                if chunk_text:
                    chunks.append(chunk_text)
                current_chunk = []
                current_size = 0

        def _is_heading(line: str) -> tuple:
            match = re.match(r'^(#{1,6})\s+(.+)', line)
            if match:
                return True, len(match.group(1))
            match = re.match(r'^(\d+(?:\.\d+)*)\s+(.+)', line)
            if match and len(match.group(1)) <= 5:
                return True, match.group(1).count(".") + 1
            return False, 0

        def _is_table_row(line: str) -> bool:
            stripped = line.strip()
            if not stripped:
                return False
            if stripped.startswith("|") and stripped.endswith("|"):
                return True
            parts = re.split(r'\s{2,}', stripped)
            if len(parts) >= 3 and any(re.search(r'\d', p) for p in parts):
                return True
            return False

        def _is_table_separator(line: str) -> bool:
            stripped = line.strip()
            return stripped.startswith("|") and "-" in stripped and set(stripped.replace("|", "").strip()) <= set("-: ")

        i = 0
        while i < len(lines):
            line = lines[i]
            line_len = len(line) + 1
            is_heading, _ = _is_heading(line)

            if is_heading and current_chunk:
                _flush_chunk()
                current_chunk = [line]
                current_size = line_len
                i += 1
                continue

            if _is_table_row(line) or _is_table_separator(line):
                table_lines = []
                while i < len(lines) and (_is_table_row(lines[i]) or _is_table_separator(lines[i])):
                    table_lines.append(lines[i])
                    i += 1
                table_text = "\n".join(table_lines)
                table_size = len(table_text)

                # Small/medium table: treat atomically
                if table_size <= max_chunk_size * 1.5:
                    if current_size + table_size > max_chunk_size and current_chunk:
                        _flush_chunk()
                    current_chunk.extend(table_lines)
                    current_size += table_size
                else:
                    # Large table: split into sub-tables with header preserved
                    if current_chunk:
                        _flush_chunk()

                    # Identify header boundary (first separator line marks header end)
                    sep_idx = None
                    for j, tl in enumerate(table_lines):
                        if _is_table_separator(tl):
                            sep_idx = j
                            break
                    header_rows = table_lines[:sep_idx + 1] if sep_idx is not None else table_lines[:2]
                    header_text = "\n".join(header_rows)
                    header_len = len(header_text) + 1
                    data_rows = table_lines[sep_idx + 1:] if sep_idx is not None else table_lines[2:]

                    sub = []
                    sub_len = header_len
                    for dr in data_rows:
                        dr_len = len(dr) + 1
                        if sub_len + dr_len > max_chunk_size and sub:
                            chunk_text = (header_text + "\n" + "\n".join(sub)).strip()
                            if chunk_text:
                                chunks.append(chunk_text)
                            sub = []
                            sub_len = header_len
                        sub.append(dr)
                        sub_len += dr_len

                    if sub:
                        chunk_text = (header_text + "\n" + "\n".join(sub)).strip()
                        if chunk_text:
                            chunks.append(chunk_text)

                    current_chunk = []
                    current_size = 0
                continue

            if current_size + line_len > max_chunk_size and current_chunk:
                _flush_chunk()

            # If a single line exceeds max_chunk_size, split it into smaller pieces
            if line_len > max_chunk_size:
                for pos in range(0, len(line), max_chunk_size):
                    piece = line[pos:pos + max_chunk_size]
                    if not piece.strip():
                        continue
                    if current_size + len(piece) > max_chunk_size and current_chunk:
                        _flush_chunk()
                    current_chunk.append(piece)
                    current_size += len(piece) + 1
            else:
                current_chunk.append(line)
                current_size += line_len
            i += 1

        _flush_chunk()
        return chunks if chunks else [text]

    def _simple_chunk_by_heading(self, text: str) -> list[str]:
        lines = text.split("\n")
        chunks = []
        current_chunk = []
        for line in lines:
            if line.startswith("#") and current_chunk:
                chunks.append("\n".join(current_chunk))
                current_chunk = [line]
            else:
                current_chunk.append(line)
        if current_chunk:
            chunks.append("\n".join(current_chunk))
        return chunks if chunks else [text]

    # ========================================================================
    # Section Title Extraction
    # ========================================================================

    _NON_CONTENT_TITLES = [
        'table of content', 'contents', 'figure index', 'table index',
        'warranty disclaimer', 'declaration', 'confidential',
        'revision history', 'about this document', 'reference',
        'list of figures', 'list of tables', 'abbreviations',
    ]

    def _is_non_content_title(self, title: str) -> bool:
        import re
        if not title:
            return True
        t_lower = title.lower().strip()
        if t_lower.startswith('_'):
            return True
        if t_lower.isdigit():
            return True
        if re.match(r'^(\d+\.)+\s*$', t_lower):
            return True
        for pattern in self._NON_CONTENT_TITLES:
            if pattern in t_lower:
                return True
        return False

    def _build_outline_map(self, pypdf_reader) -> dict:
        outline_map = {}
        try:
            outline = pypdf_reader.outline
            if not outline:
                return outline_map

            def _walk_outline(items, level=0):
                for item in items:
                    if isinstance(item, list):
                        _walk_outline(item, level + 1)
                    elif hasattr(item, 'title') and hasattr(item, 'page'):
                        title = item.title.strip() if item.title else ""
                        try:
                            page_num = pypdf_reader.get_page_number(item.page) + 1  # 1-based
                        except Exception:
                            page_num = None
                        if page_num and title and not self._is_non_content_title(title):
                            if page_num not in outline_map:
                                outline_map[page_num] = title

            _walk_outline(outline)
        except Exception as e:
            logger.warning(f"PDF outline extraction failed: {e}")
        return outline_map

    def _build_full_toc(self, pypdf_reader) -> list:
        """Extract the COMPLETE TOC from the PDF outline as an ordered entry list.

        Unlike _build_outline_map (page -> first title, used for page labeling),
        this preserves every entry including multiple sections sharing one page
        (datasheet style: 1.2.4 and 1.2.5 both start on p12). Entry order follows
        document order (outline walk), which downstream section-range building
        relies on.

        Returns: [(level, title, page_num), ...] (1-based page numbers)
        """
        toc = []
        try:
            outline = pypdf_reader.outline
            if not outline:
                return toc

            def _walk_outline(items, level=0):
                for item in items:
                    if isinstance(item, list):
                        _walk_outline(item, level + 1)
                    elif hasattr(item, 'title') and hasattr(item, 'page'):
                        title = item.title.strip() if item.title else ""
                        try:
                            page_num = pypdf_reader.get_page_number(item.page) + 1  # 1-based
                        except Exception:
                            page_num = None
                        if page_num and title and not self._is_non_content_title(title):
                            toc.append((level + 1, title, page_num))

            _walk_outline(outline)
        except Exception as e:
            logger.warning(f"PDF full TOC extraction failed: {e}")
        return toc

    def _extract_section_title(self, text: str, page_num: int,
                               outline_map: dict = None) -> str:
        import re
        if outline_map and page_num in outline_map:
            return outline_map[page_num]

        lines = text.split('\n')
        for line in lines[:5]:
            line = line.strip()
            if len(line) > 80:
                continue
            if re.match(r'^(#{1,6})\s+(.+)', line):
                return line.lstrip('#').strip()
            if re.match(r'^(\d+(?:\.\d+)*)\s+(.+)', line):
                return line
            if re.match(r'^(Chapter\s+\d+|Section\s+\d+)', line, re.IGNORECASE):
                return line
        return ""
